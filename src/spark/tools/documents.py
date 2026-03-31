"""Document tools — read Word, Excel, PowerPoint, and PDF files."""

from __future__ import annotations

from pathlib import Path
from typing import Any

_TOOLS = [
    {
        "name": "read_word",
        "description": "Extract text from a Word (.docx) document.",
        "inputSchema": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "Path to the .docx file."},
            },
            "required": ["path"],
        },
    },
    {
        "name": "read_excel",
        "description": "Read data from an Excel (.xlsx) spreadsheet.",
        "inputSchema": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "Path to the .xlsx file."},
                "sheet": {"type": "string", "description": "Sheet name. Default: first sheet."},
                "max_rows": {"type": "integer", "description": "Max rows to read. Default: 1000."},
            },
            "required": ["path"],
        },
    },
    {
        "name": "read_pdf",
        "description": "Extract text from a PDF document.",
        "inputSchema": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "Path to the PDF file."},
                "max_pages": {"type": "integer", "description": "Max pages to read. Default: 100."},
            },
            "required": ["path"],
        },
    },
    {
        "name": "read_powerpoint",
        "description": "Extract text from a PowerPoint (.pptx) presentation.",
        "inputSchema": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "Path to the .pptx file."},
            },
            "required": ["path"],
        },
    },
]


def get_tools(mode: str = "read") -> list[dict[str, Any]]:
    """Return document tool definitions."""
    return list(_TOOLS)


def execute(tool_name: str, tool_input: dict[str, Any]) -> str:
    """Execute a document tool."""
    path = Path(tool_input["path"]).resolve()
    if not path.is_file():
        return f"File not found: {path}"

    if tool_name == "read_word":
        return _read_word(path)
    elif tool_name == "read_excel":
        return _read_excel(path, tool_input.get("sheet"), tool_input.get("max_rows", 1000))
    elif tool_name == "read_pdf":
        return _read_pdf(path, tool_input.get("max_pages", 100))
    elif tool_name == "read_powerpoint":
        return _read_powerpoint(path)

    return f"Unknown document tool: {tool_name}"


def _read_word(path: Path) -> str:
    from docx import Document

    doc = Document(str(path))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    if not paragraphs:
        return "(Document is empty)"
    return "\n\n".join(paragraphs)


def _read_excel(path: Path, sheet: str | None, max_rows: int) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    ws = wb[sheet] if sheet and sheet in wb.sheetnames else wb.active
    if ws is None:
        return "No active sheet found."

    rows: list[str] = []
    for i, row in enumerate(ws.iter_rows(values_only=True)):
        if i >= max_rows:
            rows.append(f"... ({max_rows} row limit)")
            break
        cells = [str(c) if c is not None else "" for c in row]
        rows.append("\t".join(cells))

    wb.close()
    if not rows:
        return "(Sheet is empty)"
    return "\n".join(rows)


def _read_pdf(path: Path, max_pages: int) -> str:
    import pdfplumber

    pages: list[str] = []
    with pdfplumber.open(str(path)) as pdf:
        for i, page in enumerate(pdf.pages):
            if i >= max_pages:
                pages.append(f"... ({max_pages} page limit)")
                break
            text = page.extract_text()
            if text:
                pages.append(f"--- Page {i + 1} ---\n{text}")

    if not pages:
        return "(PDF has no extractable text)"
    return "\n\n".join(pages)


def _read_powerpoint(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    slides: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
        if texts:
            slides.append(f"--- Slide {i} ---\n" + "\n".join(texts))

    if not slides:
        return "(Presentation is empty)"
    return "\n\n".join(slides)
