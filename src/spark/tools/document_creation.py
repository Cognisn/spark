"""Document creation tools — create Word, Excel, PowerPoint, and PDF files.

Supports advanced formatting including headings, tables, images, charts,
styles, colours, and page layout.
"""

from __future__ import annotations

import json
import logging
from pathlib import Path
from typing import Any

logger = logging.getLogger(__name__)

TOOLS = [
    {
        "name": "create_word",
        "description": (
            "Create a Microsoft Word (.docx) document with rich formatting. "
            "Supports headings, paragraphs, tables, bullet/numbered lists, "
            "bold/italic/underline text, images, page breaks, and custom styles. "
            "Use get_tool_documentation('create_word') for full formatting guide."
        ),
        "inputSchema": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "Output file path (.docx)."},
                "content": {
                    "type": "array",
                    "description": "Array of content blocks in order.",
                    "items": {
                        "type": "object",
                        "properties": {
                            "type": {
                                "type": "string",
                                "enum": [
                                    "heading",
                                    "paragraph",
                                    "table",
                                    "list",
                                    "image",
                                    "page_break",
                                ],
                            },
                            "text": {"type": "string"},
                            "level": {"type": "integer", "description": "Heading level 1-4."},
                            "bold": {"type": "boolean"},
                            "italic": {"type": "boolean"},
                            "underline": {"type": "boolean"},
                            "colour": {
                                "type": "string",
                                "description": "Hex colour e.g. '#FF0000'.",
                            },
                            "font_size": {
                                "type": "integer",
                                "description": "Font size in points.",
                            },
                            "alignment": {
                                "type": "string",
                                "enum": ["left", "center", "right", "justify"],
                            },
                            "rows": {
                                "type": "array",
                                "description": "Table rows as arrays of cell strings.",
                                "items": {
                                    "type": "array",
                                    "items": {"type": "string"},
                                },
                            },
                            "header_row": {
                                "type": "boolean",
                                "description": "Style first row as header.",
                            },
                            "items": {
                                "type": "array",
                                "description": "List items.",
                                "items": {"type": "string"},
                            },
                            "ordered": {
                                "type": "boolean",
                                "description": "Numbered list if true.",
                            },
                            "image_path": {"type": "string"},
                            "width_inches": {"type": "number"},
                        },
                        "required": ["type"],
                    },
                },
                "title": {"type": "string", "description": "Document title metadata."},
                "author": {"type": "string", "description": "Document author metadata."},
            },
            "required": ["path", "content"],
        },
    },
    {
        "name": "create_excel",
        "description": (
            "Create a Microsoft Excel (.xlsx) workbook with formatted sheets. "
            "Supports multiple sheets, column headers with auto-filter, "
            "bold/colour formatting, column widths, number formats, and formulas. "
            "Use get_tool_documentation('create_excel') for full formatting guide."
        ),
        "inputSchema": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "Output file path (.xlsx)."},
                "sheets": {
                    "type": "array",
                    "description": "Array of sheet definitions.",
                    "items": {
                        "type": "object",
                        "properties": {
                            "name": {"type": "string", "description": "Sheet name."},
                            "headers": {
                                "type": "array",
                                "items": {"type": "string"},
                                "description": "Column header labels.",
                            },
                            "rows": {
                                "type": "array",
                                "description": "Data rows.",
                                "items": {
                                    "type": "array",
                                    "items": {},
                                },
                            },
                            "column_widths": {
                                "type": "array",
                                "items": {"type": "number"},
                                "description": "Column widths in characters.",
                            },
                            "header_colour": {
                                "type": "string",
                                "description": "Header background hex colour.",
                            },
                            "auto_filter": {
                                "type": "boolean",
                                "description": "Enable auto-filter on headers.",
                            },
                            "freeze_panes": {
                                "type": "string",
                                "description": "Cell ref to freeze at e.g. 'A2'.",
                            },
                        },
                        "required": ["name", "rows"],
                    },
                },
                "title": {"type": "string", "description": "Workbook title metadata."},
                "author": {"type": "string", "description": "Workbook author metadata."},
            },
            "required": ["path", "sheets"],
        },
    },
    {
        "name": "create_powerpoint",
        "description": (
            "Create a Microsoft PowerPoint (.pptx) presentation. "
            "Supports title slides, content slides with bullet points, "
            "tables, images, two-column layouts, and custom colours. "
            "Use get_tool_documentation('create_powerpoint') for full formatting guide."
        ),
        "inputSchema": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "Output file path (.pptx)."},
                "slides": {
                    "type": "array",
                    "description": "Array of slide definitions.",
                    "items": {
                        "type": "object",
                        "properties": {
                            "layout": {
                                "type": "string",
                                "enum": [
                                    "title",
                                    "content",
                                    "two_column",
                                    "table",
                                    "image",
                                    "blank",
                                ],
                                "description": "Slide layout type.",
                            },
                            "title": {"type": "string"},
                            "subtitle": {"type": "string"},
                            "bullets": {
                                "type": "array",
                                "items": {"type": "string"},
                            },
                            "left_bullets": {
                                "type": "array",
                                "items": {"type": "string"},
                            },
                            "right_bullets": {
                                "type": "array",
                                "items": {"type": "string"},
                            },
                            "rows": {
                                "type": "array",
                                "items": {
                                    "type": "array",
                                    "items": {"type": "string"},
                                },
                            },
                            "image_path": {"type": "string"},
                            "notes": {"type": "string"},
                        },
                        "required": ["layout"],
                    },
                },
                "title": {"type": "string", "description": "Presentation title metadata."},
                "author": {"type": "string", "description": "Presentation author metadata."},
            },
            "required": ["path", "slides"],
        },
    },
    {
        "name": "create_pdf",
        "description": (
            "Create a PDF document with formatted content. "
            "Supports headings, paragraphs, tables, bullet lists, images, "
            "page breaks, headers/footers, and custom fonts/colours. "
            "Use get_tool_documentation('create_pdf') for full formatting guide."
        ),
        "inputSchema": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "Output file path (.pdf)."},
                "content": {
                    "type": "array",
                    "description": "Array of content blocks in order.",
                    "items": {
                        "type": "object",
                        "properties": {
                            "type": {
                                "type": "string",
                                "enum": [
                                    "heading",
                                    "paragraph",
                                    "table",
                                    "list",
                                    "image",
                                    "page_break",
                                    "spacer",
                                ],
                            },
                            "text": {"type": "string"},
                            "level": {"type": "integer"},
                            "bold": {"type": "boolean"},
                            "italic": {"type": "boolean"},
                            "colour": {"type": "string"},
                            "font_size": {"type": "integer"},
                            "alignment": {
                                "type": "string",
                                "enum": ["left", "center", "right"],
                            },
                            "rows": {
                                "type": "array",
                                "items": {
                                    "type": "array",
                                    "items": {"type": "string"},
                                },
                            },
                            "items": {
                                "type": "array",
                                "items": {"type": "string"},
                            },
                            "image_path": {"type": "string"},
                            "width": {"type": "number", "description": "Image width in inches."},
                            "height": {"type": "number", "description": "Spacer height in points."},
                        },
                        "required": ["type"],
                    },
                },
                "title": {"type": "string", "description": "Document title (shown in header)."},
                "author": {"type": "string", "description": "Document author metadata."},
                "page_size": {
                    "type": "string",
                    "enum": ["A4", "letter"],
                    "description": "Page size. Default: A4.",
                },
            },
            "required": ["path", "content"],
        },
    },
]


def get_tools() -> list[dict[str, Any]]:
    """Return document creation tool definitions."""
    return list(TOOLS)


def execute(tool_name: str, tool_input: dict[str, Any], allowed_paths: list[str]) -> str:
    """Execute a document creation tool."""
    path = Path(tool_input["path"]).resolve()

    # Validate against allowed paths
    if allowed_paths:
        if not any(str(path).startswith(str(Path(ap).resolve())) for ap in allowed_paths):
            return f"Access denied: {path} is outside allowed paths."

    # Ensure parent directory exists
    path.parent.mkdir(parents=True, exist_ok=True)

    try:
        if tool_name == "create_word":
            return _create_word(path, tool_input)
        elif tool_name == "create_excel":
            return _create_excel(path, tool_input)
        elif tool_name == "create_powerpoint":
            return _create_powerpoint(path, tool_input)
        elif tool_name == "create_pdf":
            return _create_pdf(path, tool_input)
        return f"Unknown document creation tool: {tool_name}"
    except Exception as e:
        logger.error("Document creation error (%s): %s", tool_name, e, exc_info=True)
        return f"Error creating document: {e}"


# ---------------------------------------------------------------------------
# Word (.docx)
# ---------------------------------------------------------------------------


def _create_word(path: Path, tool_input: dict) -> str:
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Inches, Pt, RGBColor

    doc = Document()

    # Metadata
    if tool_input.get("title"):
        doc.core_properties.title = tool_input["title"]
    if tool_input.get("author"):
        doc.core_properties.author = tool_input["author"]

    align_map = {
        "left": WD_ALIGN_PARAGRAPH.LEFT,
        "center": WD_ALIGN_PARAGRAPH.CENTER,
        "right": WD_ALIGN_PARAGRAPH.RIGHT,
        "justify": WD_ALIGN_PARAGRAPH.JUSTIFY,
    }

    for block in tool_input.get("content", []):
        block_type = block.get("type", "")

        if block_type == "heading":
            level = min(max(block.get("level", 1), 1), 4)
            p = doc.add_heading(block.get("text", ""), level=level)
            if block.get("alignment"):
                p.alignment = align_map.get(block["alignment"])

        elif block_type == "paragraph":
            p = doc.add_paragraph()
            if block.get("alignment"):
                p.alignment = align_map.get(block["alignment"])
            run = p.add_run(block.get("text", ""))
            if block.get("bold"):
                run.bold = True
            if block.get("italic"):
                run.italic = True
            if block.get("underline"):
                run.underline = True
            if block.get("font_size"):
                run.font.size = Pt(block["font_size"])
            if block.get("colour"):
                try:
                    hex_colour = block["colour"].lstrip("#")
                    run.font.color.rgb = RGBColor(
                        int(hex_colour[0:2], 16),
                        int(hex_colour[2:4], 16),
                        int(hex_colour[4:6], 16),
                    )
                except (ValueError, IndexError):
                    pass

        elif block_type == "table":
            rows = block.get("rows", [])
            if rows:
                table = doc.add_table(rows=len(rows), cols=len(rows[0]))
                table.style = "Table Grid"
                for i, row_data in enumerate(rows):
                    for j, cell_text in enumerate(row_data):
                        table.rows[i].cells[j].text = str(cell_text)
                # Bold header row
                if block.get("header_row", True) and rows:
                    for cell in table.rows[0].cells:
                        for paragraph in cell.paragraphs:
                            for run in paragraph.runs:
                                run.bold = True

        elif block_type == "list":
            items = block.get("items", [])
            ordered = block.get("ordered", False)
            for i, item in enumerate(items):
                style = "List Number" if ordered else "List Bullet"
                doc.add_paragraph(item, style=style)

        elif block_type == "image":
            img_path = block.get("image_path", "")
            if img_path and Path(img_path).is_file():
                width = Inches(block.get("width_inches", 5))
                doc.add_picture(img_path, width=width)

        elif block_type == "page_break":
            doc.add_page_break()

    doc.save(str(path))
    return f"Word document created: {path}"


# ---------------------------------------------------------------------------
# Excel (.xlsx)
# ---------------------------------------------------------------------------


def _create_excel(path: Path, tool_input: dict) -> str:
    from openpyxl import Workbook
    from openpyxl.styles import Alignment, Font, PatternFill
    from openpyxl.utils import get_column_letter

    wb = Workbook()
    # Remove default sheet
    if wb.active:
        wb.remove(wb.active)

    # Metadata
    if tool_input.get("title"):
        wb.properties.title = tool_input["title"]
    if tool_input.get("author"):
        wb.properties.creator = tool_input["author"]

    for sheet_def in tool_input.get("sheets", []):
        ws = wb.create_sheet(title=sheet_def.get("name", "Sheet"))

        headers = sheet_def.get("headers", [])
        rows = sheet_def.get("rows", [])
        header_colour = sheet_def.get("header_colour", "4472C4")
        if header_colour.startswith("#"):
            header_colour = header_colour[1:]

        # Write headers
        if headers:
            header_fill = PatternFill(start_color=header_colour, fill_type="solid")
            header_font = Font(bold=True, color="FFFFFF")
            for col_idx, header in enumerate(headers, 1):
                cell = ws.cell(row=1, column=col_idx, value=header)
                cell.fill = header_fill
                cell.font = header_font
                cell.alignment = Alignment(horizontal="center")

        # Write data rows
        start_row = 2 if headers else 1
        for row_idx, row_data in enumerate(rows, start_row):
            for col_idx, value in enumerate(row_data, 1):
                cell = ws.cell(row=row_idx, column=col_idx)
                # Try to parse numbers
                if isinstance(value, str):
                    try:
                        cell.value = float(value) if "." in value else int(value)
                    except (ValueError, TypeError):
                        cell.value = value
                else:
                    cell.value = value

        # Column widths
        col_widths = sheet_def.get("column_widths", [])
        if col_widths:
            for i, width in enumerate(col_widths, 1):
                ws.column_dimensions[get_column_letter(i)].width = width
        else:
            # Auto-size based on content
            max_cols = len(headers) or (len(rows[0]) if rows else 0)
            for col_idx in range(1, max_cols + 1):
                max_len = 10
                for row in ws.iter_rows(min_col=col_idx, max_col=col_idx, values_only=False):
                    for cell in row:
                        if cell.value:
                            max_len = max(max_len, len(str(cell.value)) + 2)
                ws.column_dimensions[get_column_letter(col_idx)].width = min(max_len, 50)

        # Auto-filter
        if sheet_def.get("auto_filter", True) and headers:
            last_col = get_column_letter(len(headers))
            ws.auto_filter.ref = f"A1:{last_col}1"

        # Freeze panes
        freeze = sheet_def.get("freeze_panes", "A2" if headers else None)
        if freeze:
            ws.freeze_panes = freeze

    wb.save(str(path))
    total_rows = sum(len(s.get("rows", [])) for s in tool_input.get("sheets", []))
    return f"Excel workbook created: {path} ({len(tool_input.get('sheets', []))} sheets, {total_rows} data rows)"


# ---------------------------------------------------------------------------
# PowerPoint (.pptx)
# ---------------------------------------------------------------------------


def _create_powerpoint(path: Path, tool_input: dict) -> str:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()

    # Metadata
    if tool_input.get("title"):
        prs.core_properties.title = tool_input["title"]
    if tool_input.get("author"):
        prs.core_properties.author = tool_input["author"]

    for slide_def in tool_input.get("slides", []):
        layout_name = slide_def.get("layout", "content")

        if layout_name == "title":
            slide_layout = prs.slide_layouts[0]  # Title Slide
            slide = prs.slides.add_slide(slide_layout)
            if slide.shapes.title:
                slide.shapes.title.text = slide_def.get("title", "")
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = slide_def.get("subtitle", "")

        elif layout_name == "content":
            slide_layout = prs.slide_layouts[1]  # Title and Content
            slide = prs.slides.add_slide(slide_layout)
            if slide.shapes.title:
                slide.shapes.title.text = slide_def.get("title", "")
            if len(slide.placeholders) > 1:
                tf = slide.placeholders[1].text_frame
                tf.clear()
                for i, bullet in enumerate(slide_def.get("bullets", [])):
                    if i == 0:
                        tf.paragraphs[0].text = bullet
                    else:
                        p = tf.add_paragraph()
                        p.text = bullet

        elif layout_name == "two_column":
            slide_layout = prs.slide_layouts[3]  # Two Content
            slide = prs.slides.add_slide(slide_layout)
            if slide.shapes.title:
                slide.shapes.title.text = slide_def.get("title", "")
            # Left column
            if len(slide.placeholders) > 1:
                tf = slide.placeholders[1].text_frame
                tf.clear()
                for i, item in enumerate(slide_def.get("left_bullets", [])):
                    if i == 0:
                        tf.paragraphs[0].text = item
                    else:
                        tf.add_paragraph().text = item
            # Right column
            if len(slide.placeholders) > 2:
                tf = slide.placeholders[2].text_frame
                tf.clear()
                for i, item in enumerate(slide_def.get("right_bullets", [])):
                    if i == 0:
                        tf.paragraphs[0].text = item
                    else:
                        tf.add_paragraph().text = item

        elif layout_name == "table":
            slide_layout = prs.slide_layouts[5]  # Blank
            slide = prs.slides.add_slide(slide_layout)
            # Add title as text box
            if slide_def.get("title"):
                from pptx.util import Emu

                txbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
                tf = txbox.text_frame
                tf.text = slide_def["title"]
                tf.paragraphs[0].font.size = Pt(24)
                tf.paragraphs[0].font.bold = True

            rows_data = slide_def.get("rows", [])
            if rows_data:
                num_rows = len(rows_data)
                num_cols = max(len(r) for r in rows_data)
                table = slide.shapes.add_table(
                    num_rows, num_cols, Inches(0.5), Inches(1.5), Inches(9), Inches(0.4 * num_rows)
                ).table
                for i, row_data in enumerate(rows_data):
                    for j, cell_text in enumerate(row_data):
                        table.cell(i, j).text = str(cell_text)

        elif layout_name == "image":
            slide_layout = prs.slide_layouts[5]  # Blank
            slide = prs.slides.add_slide(slide_layout)
            if slide_def.get("title"):
                from pptx.util import Emu

                txbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
                txbox.text_frame.text = slide_def["title"]
                txbox.text_frame.paragraphs[0].font.size = Pt(24)
                txbox.text_frame.paragraphs[0].font.bold = True
            img_path = slide_def.get("image_path", "")
            if img_path and Path(img_path).is_file():
                slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), Inches(8))

        elif layout_name == "blank":
            slide_layout = prs.slide_layouts[6]  # Blank
            slide = prs.slides.add_slide(slide_layout)

        # Speaker notes
        if slide_def.get("notes"):
            slide.notes_slide.notes_text_frame.text = slide_def["notes"]

    prs.save(str(path))
    return f"PowerPoint presentation created: {path} ({len(tool_input.get('slides', []))} slides)"


# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------


def _create_pdf(path: Path, tool_input: dict) -> str:
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4, letter
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import inch
    from reportlab.platypus import (
        Image,
        PageBreak,
        Paragraph,
        SimpleDocTemplate,
        Spacer,
        Table,
        TableStyle,
    )

    page_size = A4 if tool_input.get("page_size", "A4") == "A4" else letter
    doc_title = tool_input.get("title", "")
    doc_author = tool_input.get("author", "")

    doc = SimpleDocTemplate(
        str(path),
        pagesize=page_size,
        title=doc_title,
        author=doc_author,
        topMargin=0.75 * inch,
        bottomMargin=0.75 * inch,
        leftMargin=0.75 * inch,
        rightMargin=0.75 * inch,
    )

    styles = getSampleStyleSheet()

    # Custom styles for different heading levels
    heading_styles = {
        1: ParagraphStyle("H1", parent=styles["Heading1"], fontSize=18, spaceAfter=12),
        2: ParagraphStyle("H2", parent=styles["Heading2"], fontSize=15, spaceAfter=10),
        3: ParagraphStyle("H3", parent=styles["Heading3"], fontSize=12, spaceAfter=8),
        4: ParagraphStyle("H4", parent=styles["Heading4"], fontSize=11, spaceAfter=6),
    }

    align_map = {"left": 0, "center": 1, "right": 2}

    flowables = []

    for block in tool_input.get("content", []):
        block_type = block.get("type", "")

        if block_type == "heading":
            level = min(max(block.get("level", 1), 1), 4)
            style = heading_styles.get(level, heading_styles[1])
            text = block.get("text", "")
            if block.get("colour"):
                text = f'<font color="{block["colour"]}">{text}</font>'
            flowables.append(Paragraph(text, style))

        elif block_type == "paragraph":
            style = ParagraphStyle("custom", parent=styles["Normal"], fontSize=10)
            if block.get("font_size"):
                style.fontSize = block["font_size"]
            if block.get("alignment"):
                style.alignment = align_map.get(block["alignment"], 0)

            text = block.get("text", "")
            if block.get("bold"):
                text = f"<b>{text}</b>"
            if block.get("italic"):
                text = f"<i>{text}</i>"
            if block.get("colour"):
                text = f'<font color="{block["colour"]}">{text}</font>'
            flowables.append(Paragraph(text, style))

        elif block_type == "table":
            rows = block.get("rows", [])
            if rows:
                table_data = [[str(c) for c in row] for row in rows]
                t = Table(table_data, repeatRows=1)
                style_commands = [
                    ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#4472C4")),
                    ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                    ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                    ("FONTSIZE", (0, 0), (-1, -1), 9),
                    ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
                    (
                        "ROWBACKGROUNDS",
                        (0, 1),
                        (-1, -1),
                        [colors.white, colors.HexColor("#F2F2F2")],
                    ),
                    ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                    ("TOPPADDING", (0, 0), (-1, -1), 4),
                    ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
                ]
                t.setStyle(TableStyle(style_commands))
                flowables.append(t)

        elif block_type == "list":
            items = block.get("items", [])
            for i, item in enumerate(items):
                prefix = f"{i + 1}. " if block.get("ordered") else "\u2022 "
                bullet_style = ParagraphStyle(
                    "bullet",
                    parent=styles["Normal"],
                    fontSize=10,
                    leftIndent=20,
                    firstLineIndent=-12,
                )
                flowables.append(Paragraph(f"{prefix}{item}", bullet_style))

        elif block_type == "image":
            img_path = block.get("image_path", "")
            if img_path and Path(img_path).is_file():
                width = block.get("width", 5) * inch
                flowables.append(Image(img_path, width=width))

        elif block_type == "page_break":
            flowables.append(PageBreak())

        elif block_type == "spacer":
            height = block.get("height", 12)
            flowables.append(Spacer(1, height))

    if not flowables:
        flowables.append(Paragraph("(Empty document)", styles["Normal"]))

    doc.build(flowables)
    return f"PDF document created: {path}"
